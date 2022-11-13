from pptx import Presentation
import os
import pandas as pd
import argparse
import time


def main():
    parser = argparse.ArgumentParser(description='Create certificate from template')
    parser.add_argument('-t', '--template', type=str, help='.pptx template file path', required=True)
    parser.add_argument('-l', '--list', type=str, help='.csv list path', required=True)
    parser.add_argument('-o', '--output', type=str, help='output dir', default='output')
    parser.add_argument('-p', '--pptx', type=bool, help='output pptx', default=False)
    args = parser.parse_args()

    if not os.path.exists(args.template):
        print('Template file not found')
        exit(1)

    if not os.path.exists(args.list):
        print('List file not found')
        exit(1)

    out_dir = args.output + '-' + time.strftime("%Y%m%d-%H%M%S")
    if not os.path.exists(out_dir):
        os.mkdir(out_dir)

    df = pd.read_csv(args.list)

    for index, row in df.iterrows():
        print(row.to_dict())
        create_certificate_from_template(args.template, os.path.join(out_dir,
                                         str(row['name'])+'.pptx'), row.to_dict(),
                                         )


def create_certificate_from_template(template, target, change_dict):
    ppt = Presentation(template)

    for s, slide in enumerate(ppt.slides):
        for ss, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p, paragraph in enumerate(shape.text_frame.paragraphs):
                for r, run in enumerate(paragraph.runs):
                    for phrase in change_dict:
                        if run.text == phrase:
                            ppt.slides[s].shapes[ss].text_frame.paragraphs[p].runs[r].text = change_dict[phrase]
    ppt.save(target)

    # for Linux
    os.system(f'unoconv -f pdf "{target}"')
    os.system(f'convert -density 300 {target[:-5]}.pdf -quality 100 {target[:-5]}.jpg')
    os.system(f'rm "{os.path.abspath(target)}"')
    # for windows
    # os.system('unoconv.bat -f pdf "{}"'.format(os.path.abspath(target)))

    print("Done")


if __name__ == "__main__":

    main()
